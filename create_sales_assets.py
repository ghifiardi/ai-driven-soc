#!/usr/bin/env python3
"""
Sales Asset Generator
=====================

Converts MSSP_COMMERCIAL_ARCHITECTURE.md into:
1. MSSP_Commercial_Architecture.docx (Word Document)
2. MSSP_Sales_Presentation.pptx (PowerPoint Deck)

Requires:
    pip install python-docx python-pptx markdown
"""

import os
import re
from docx import Document
from docx.shared import Pt, Inches, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from pptx.enum.text import PP_ALIGN

SOURCE_FILE = "MSSP_COMMERCIAL_ARCHITECTURE.md"
DOCX_OUTPUT = "MSSP_Commercial_Architecture.docx"
PPTX_OUTPUT = "MSSP_Sales_Presentation.pptx"

def read_markdown(filepath):
    with open(filepath, 'r') as f:
        return f.read()

def parse_markdown(content):
    """Simple parser to extract headers and content."""
    slides = []
    current_slide = {"title": "Introduction", "content": []}
    
    lines = content.split('\n')
    for line in lines:
        if line.startswith('# '):
            # Main Title
            current_slide["title"] = line[2:].strip()
        elif line.startswith('## '):
            # New Slide/Section
            if current_slide["content"]:
                slides.append(current_slide)
            current_slide = {"title": line[3:].strip(), "content": []}
        elif line.startswith('### '):
            # Sub-header
            current_slide["content"].append(f"BOLD:{line[4:].strip()}")
        elif line.startswith('* ') or line.startswith('- '):
            # Bullet point
            current_slide["content"].append(f"BULLET:{line[2:].strip()}")
        elif line.strip() and not line.startswith('```'):
            # Normal text
            current_slide["content"].append(f"TEXT:{line.strip()}")
            
    if current_slide["content"]:
        slides.append(current_slide)
        
    return slides

def create_docx(slides):
    doc = Document()
    
    # Title Style
    style = doc.styles['Normal']
    font = style.font
    font.name = 'Calibri'
    font.size = Pt(11)
    
    # Add Title
    title = doc.add_heading('AI-Driven SOC Platform', 0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    doc.add_paragraph('Commercial Architecture & Value Proposition', style='Subtitle')
    doc.add_paragraph('')
    
    for slide in slides:
        # Section Header
        h = doc.add_heading(slide['title'], level=1)
        
        # Check for Diagram Section
        if "Technical Architecture" in slide['title'] and os.path.exists("architecture_diagram.png"):
            doc.add_picture("architecture_diagram.png", width=Inches(6))
        
        for item in slide['content']:
            if item.startswith("BOLD:"):
                p = doc.add_paragraph()
                runner = p.add_run(item[5:])
                runner.bold = True
            elif item.startswith("BULLET:"):
                doc.add_paragraph(item[7:], style='List Bullet')
            elif item.startswith("TEXT:"):
                if not item.startswith("TEXT:```"): # Skip code blocks if we have the image
                    doc.add_paragraph(item[5:])
                
    doc.save(DOCX_OUTPUT)
    print(f"Generated {DOCX_OUTPUT}")

def create_pptx(slides):
    prs = Presentation()
    
    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "AI-Driven SOC Platform"
    subtitle.text = "Commercial Architecture & Value Proposition"
    
    # Content Slides
    bullet_slide_layout = prs.slide_layouts[1]
    blank_slide_layout = prs.slide_layouts[6]
    
    for slide_data in slides:
        if "Technical Architecture" in slide_data['title'] and os.path.exists("architecture_diagram.png"):
            slide = prs.slides.add_slide(blank_slide_layout)
            shapes = slide.shapes
            title_shape = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(0.2), PptxInches(9), PptxInches(1))
            title_shape.text_frame.text = slide_data['title']
            
            # Add Image
            slide.shapes.add_picture("architecture_diagram.png", PptxInches(1), PptxInches(1.5), width=PptxInches(8))
            continue

        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        
        title_shape = shapes.title
        title_shape.text = slide_data['title']
        
        body_shape = shapes.placeholders[1]
        tf = body_shape.text_frame
        
        for item in slide_data['content']:
            if item.startswith("BOLD:"):
                p = tf.add_paragraph()
                p.text = item[5:]
                p.font.bold = True
                p.level = 0
            elif item.startswith("BULLET:"):
                p = tf.add_paragraph()
                p.text = item[7:]
                p.level = 1
            elif item.startswith("TEXT:"):
                if len(item) < 200 and not item.startswith("TEXT:```"): 
                    p = tf.add_paragraph()
                    p.text = item[5:]
                    p.level = 0
                    
    prs.save(PPTX_OUTPUT)
    print(f"Generated {PPTX_OUTPUT}")

def main():
    if not os.path.exists(SOURCE_FILE):
        print(f"Error: {SOURCE_FILE} not found.")
        return

    print("Parsing markdown...")
    content = read_markdown(SOURCE_FILE)
    slides = parse_markdown(content)
    
    print("Creating Word Document...")
    create_docx(slides)
    
    print("Creating PowerPoint Presentation...")
    create_pptx(slides)
    
    print("Done!")

if __name__ == "__main__":
    main()
